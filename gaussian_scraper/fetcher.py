from __future__ import annotations

import time
from io import BytesIO

import requests
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

# How long to wait for a page to respond before giving up.
REQUEST_TIMEOUT_SEC = 10

# Seconds to wait between fetching pages to avoid rate limiting.
REQUEST_DELAY_SEC = 2

# Realistic browser User-Agent to avoid bot detection on documentation sites.
USER_AGENT = (
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
    "AppleWebKit/537.36 (KHTML, like Gecko) "
    "Chrome/124.0.0.0 Safari/537.36"
)

# Tags that contain visible documentation content.
# skip nav, header, footer, and script tags as they contain UI noise.
CONTENT_TAGS = ["p", "li", "pre", "code", "h1", "h2", "h3", "h4"]


def _extract_pdf_text(pdf_bytes: bytes) -> str | None:
    """
    Extract text from a PDF's raw bytes, one line per extracted text
    fragment -- mirrors fetch_page_text's HTML behavior of one line per
    content piece, so downstream keyword-density scoring works the same
    way regardless of where the text came from.

    Returns None if extraction fails or yields nothing -- e.g. a
    scanned/image-only PDF with no embedded text layer, the PDF
    equivalent of a JS-rendered page BeautifulSoup can't see into.

    Naive extraction like this can produce out-of-reading-order text for
    multi-column layouts and repeated page headers/footers. Not solved
    here -- extract_relevant_passages's per-line keyword scoring still
    surfaces genuinely relevant lines regardless of overall document
    coherence, the same way it already tolerates messy HTML extraction.
    """
    try:
        reader = PdfReader(BytesIO(pdf_bytes))
        text = "\n".join(page.extract_text() for page in reader.pages)
    except Exception:
        return None
    return text if text.strip() else None


def _extract_docx_text(docx_bytes: bytes) -> str | None:
    """
    Extract paragraph and table-cell text from a .docx's raw bytes. Same
    "one line per fragment" contract and the same accepted extraction
    limitations as _extract_pdf_text -- flattening a table to cell-by-
    cell text loses its row/column structure, but per-line keyword
    scoring downstream doesn't need that structure to find relevant text.
    """
    try:
        document = Document(BytesIO(docx_bytes))
        paragraphs = [p.text for p in document.paragraphs]
        table_cells = [
            cell.text
            for table in document.tables
            for row in table.rows
            for cell in row.cells
        ]
        text = "\n".join(paragraphs + table_cells)
    except Exception:
        return None
    return text if text.strip() else None


def _extract_pptx_text(pptx_bytes: bytes) -> str | None:
    """
    Extract text from every text-bearing shape across all slides in a
    .pptx. Speaker notes are deliberately not included -- slide content
    is the analog of a page's visible text, notes are more like a
    separate document.
    """
    try:
        presentation = Presentation(BytesIO(pptx_bytes))
        lines = [
            shape.text_frame.text
            for slide in presentation.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        text = "\n".join(lines)
    except Exception:
        return None
    return text if text.strip() else None


def fetch_page_text(url: str) -> tuple[str | None, int | None]:
    """
    Fetch a URL and return its visible text content as a plain string,
    along with the HTTP status code for logging purposes. Handles HTML
    pages, PDF, DOCX, PPTX, and plain text transparently -- callers don't
    need to know which one a URL points to.

    Steps:
    - Wait briefly between requests to avoid rate limiting
    - Make an HTTP GET request with a realistic browser User-Agent header
    - Detect the response's actual format (by Content-Type or URL suffix)
      and extract text the way that format needs: PDF/DOCX/PPTX are
      parsed from raw bytes, plain text is used as-is, legacy binary
      Office formats (.doc/.ppt/.xls) are explicitly rejected as
      unextractable (not OOXML, not HTML -- parsing them as HTML would
      produce garbled text that can spuriously match keywords), and
      anything else falls through to BeautifulSoup HTML parsing of
      content-bearing tags
    - Return a single string with one line per extracted piece

    Returns (None, None) if the request fails due to a network error.
    Returns (None, status_code) if the server responds with a non-200 status.
    Returns (text, 200) on success.
    """
    time.sleep(REQUEST_DELAY_SEC)

    try:
        response = requests.get(
            url,
            timeout=REQUEST_TIMEOUT_SEC,
            headers={"User-Agent": USER_AGENT},
        )
    except requests.RequestException:
        return None, None

    if response.status_code != 200:
        return None, response.status_code

    content_type = response.headers.get("Content-Type", "")
    lower_url = url.lower()

    if "application/pdf" in content_type or lower_url.endswith(".pdf"):
        return _extract_pdf_text(response.content), 200

    if "wordprocessingml.document" in content_type or lower_url.endswith(".docx"):
        return _extract_docx_text(response.content), 200

    if "presentationml.presentation" in content_type or lower_url.endswith(".pptx"):
        return _extract_pptx_text(response.content), 200

    if "text/plain" in content_type or lower_url.endswith(".txt"):
        text = response.text
        return (text if text.strip() else None), 200

    # Legacy pre-2007 binary Office formats (.doc/.ppt/.xls) aren't OOXML,
    # so they don't match the docx/pptx branches above and aren't valid
    # HTML either -- without this check they'd silently fall through to
    # BeautifulSoup, which parses the binary bytes as if they were HTML
    # text and can produce garbled "content" that coincidentally contains
    # keyword substrings, scoring as a false GOOD instead of correctly
    # coming back as unextractable (same as an image-only PDF).
    legacy_office_types = ("application/msword", "application/vnd.ms-powerpoint", "application/vnd.ms-excel")
    legacy_office_suffixes = (".doc", ".ppt", ".xls")
    if any(t in content_type for t in legacy_office_types) or lower_url.endswith(legacy_office_suffixes):
        return None, 200

    soup = BeautifulSoup(response.text, "html.parser")

    # Remove script and style tags so their text doesn't pollute the output.
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    lines = []
    for tag in soup.find_all(CONTENT_TAGS):
        text = tag.get_text(separator=" ", strip=True)
        if text:
            lines.append(text)

    return ("\n".join(lines) if lines else None), 200
